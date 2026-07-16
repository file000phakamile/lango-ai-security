#!/usr/bin/env python3
"""Generates pitch-deck/Lango_Pitch_Deck.pptx from the same content already
verified in pitch-deck/index.html, rather than hand-building ten slides of
shapes that could drift from it.

Kept as a committed, re-runnable script rather than shipping only the
binary .pptx — the same "a binary artifact that can silently go stale is
worse than a versionable source" reasoning that led to building the pitch
deck as HTML in the first place (see Questions.md). If the content in
docs/PITCH_DECK_CONTENT.md or the HTML deck's design ever changes, this
script is what regenerates the .pptx to match, instead of someone hand-
editing slides in PowerPoint and drifting from the source of truth.

Text slides (1, 2, 3, 5, 8, 9, 10) are built as real, editable PowerPoint
text — not images — using the same colors/fonts/proportions as the HTML
deck. Slides 4, 6, and 7 (the pipeline flow, the fairness bar comparison,
and the PSI spike chart) are inserted as full-slide background images,
captured directly from the real, already-verified index.html via a
headless browser (see capture_slides.mjs's own comments, run separately
during this task — not re-run by this script, since the PNGs it produced
are the ones actually embedded below). This guarantees those three visuals
are pixel-identical to the verified HTML deck, not a reimplementation.

Design tokens below are copied directly from the real source, confirmed by
grep immediately before writing this script, not recalled from memory:
  - Every hex color: `grep -rhoE "#[0-9A-Fa-f]{6}" components/lango/*.tsx`
  - Fonts: app/layout.tsx's next/font/google setup (IBM Plex Sans 400/500/
    600/700, IBM Plex Mono 400/500/600)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import os

# ---------------------------------------------------------------------------
# Design tokens — copied from the real dashboard source, not invented.
# ---------------------------------------------------------------------------

BG = RGBColor(0xF6, 0xF7, 0xF8)
CARD = RGBColor(0xFF, 0xFF, 0xFF)
TEXT = RGBColor(0x14, 0x17, 0x1C)
TEXT_SECONDARY = RGBColor(0x5B, 0x62, 0x70)
TEXT_MUTED = RGBColor(0x8A, 0x93, 0xA1)
BORDER = RGBColor(0xE1, 0xE4, 0xE8)
GOLD = RGBColor(0x8A, 0x63, 0x23)
GOLD_TINT = RGBColor(0xF3, 0xEC, 0xE1)  # a flat approximation of the HTML deck's
# rgba gold-at-8%-opacity tint (#8A632314) — pptx solid fills don't do alpha
# blending against an arbitrary background the way CSS rgba does, so this is
# the closest flat color, not a literal value from the source. See Questions.md.
RED = RGBColor(0xA8, 0x3A, 0x3A)
GREEN = RGBColor(0x2F, 0x7A, 0x53)

FONT_SANS = "IBM Plex Sans"
FONT_MONO = "IBM Plex Mono"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN_X = Inches(13.333 * 0.07)   # 7%, matching index.html's slide padding
MARGIN_Y = Inches(7.5 * 0.055)     # 5.5%, matching index.html's slide padding

HERE = os.path.dirname(os.path.abspath(__file__))
CAPTURES_DIR = os.environ.get("LANGO_DECK_CAPTURES", HERE)


def set_background(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    return box, tf


def add_eyebrow(slide, text):
    box, tf = add_textbox(slide, MARGIN_X, MARGIN_Y, Inches(11), Inches(0.35))
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text.upper()
    run.font.name = FONT_MONO
    run.font.size = Pt(13)
    run.font.color.rgb = GOLD
    run.font.bold = False
    # letter-spacing isn't directly exposed by python-pptx; left as a
    # documented simplification versus the HTML deck's tracked-out eyebrow
    # text — see Questions.md.
    return box


def add_headline(slide, text, top=Inches(1.0), size=30, width=Inches(11.3), align=PP_ALIGN.LEFT):
    box, tf = add_textbox(slide, MARGIN_X, top, width, Inches(1.3))
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = FONT_SANS
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = TEXT
    return box


def add_bullets(slide, items, top, width=Inches(11.3), size=16, gap_pt=10):
    box, tf = add_textbox(slide, MARGIN_X, top, width, Inches(7.5) - top - Inches(0.6))
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap_pt)
        p.line_spacing = 1.15
        # A small gold square "bullet" as a literal run character keeps this
        # a single real run of editable text (not a native PowerPoint bullet
        # list, whose default glyph/indent styling doesn't match the HTML
        # deck's square-gold-mark bullets) — a deliberate, documented
        # simplification; see Questions.md.
        r_mark = p.add_run()
        r_mark.text = "■  "
        r_mark.font.name = FONT_SANS
        r_mark.font.size = Pt(size - 3)
        r_mark.font.color.rgb = GOLD
        r_text = p.add_run()
        r_text.text = item
        r_text.font.name = FONT_SANS
        r_text.font.size = Pt(size)
        r_text.font.color.rgb = TEXT
    return box


def add_footer(slide):
    box, tf = add_textbox(slide, MARGIN_X, Inches(7.5) - Inches(0.45), Inches(6), Inches(0.3))
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "●  LANGO — AI DATA GUARD"
    run.font.name = FONT_MONO
    run.font.size = Pt(10)
    run.font.color.rgb = TEXT_MUTED


def new_slide(prs):
    blank_layout = prs.slide_layouts[6]  # "Blank" — no placeholder boxes, no template chrome
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide, BG)
    return slide


def text_slide(prs, eyebrow, headline, bullets, headline_size=30):
    slide = new_slide(prs)
    if eyebrow:
        add_eyebrow(slide, eyebrow)
    add_headline(slide, headline, size=headline_size)
    add_bullets(slide, bullets, top=Inches(2.35))
    add_footer(slide)
    return slide


def image_slide(prs, image_path):
    slide = new_slide(prs)
    slide.shapes.add_picture(image_path, 0, 0, width=SLIDE_W, height=SLIDE_H)
    return slide


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # ---- Slide 1 — Title (docs/PITCH_DECK_CONTENT.md, section 1) ----
    slide = new_slide(prs)
    box, tf = add_textbox(slide, MARGIN_X, Inches(1.7), Inches(11), Inches(0.4))
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "AI4I 2026 CHALLENGE — TRACK 4 (DEPLOYMENT)"
    run.font.name = FONT_MONO
    run.font.size = Pt(13)
    run.font.color.rgb = TEXT_MUTED
    add_headline(slide, "Lango — AI Data Guard", top=Inches(2.25), size=48)
    add_bullets(
        slide,
        [
            "Security and governance gateway for enterprise AI use",
            "AI4I 2026 Challenge — Track 4 (Deployment)",
            "Team Lango: Phakamile Mlala & Vanessa Moyo, NUST Bulawayo",
        ],
        top=Inches(3.7),
    )
    add_footer(slide)

    # ---- Slide 2 — Problem ----
    text_slide(
        prs,
        "02 — Problem",
        "Staff are pasting real institutional data into AI tools — with zero oversight",
        [
            "National IDs, bank details, and patient records routinely enter AI chat prompts",
            "No logging, no control, no way to prove what left the institution",
            "A live compliance and data-protection exposure today, not a hypothetical one",
        ],
    )

    # ---- Slide 3 — Who It's For ----
    text_slide(
        prs,
        "03 — Who It's For",
        "Built for the people accountable when this goes wrong",
        [
            "Primary user: frontline staff whose prompts pass through the gateway",
            "Beneficiary and payer: compliance, risk, and IT security teams",
            "Target sectors: banks, hospitals, government ministries",
        ],
    )

    # ---- Slide 4 — Solution / pipeline (screenshot, real visual) ----
    image_slide(prs, os.path.join(CAPTURES_DIR, "slide-4-capture.png"))

    # ---- Slide 5 — Live Demo (native, deliberate callout) ----
    slide = new_slide(prs)
    add_eyebrow(slide, "05 — Live Demo")
    add_headline(
        slide,
        "Switch to live demo here",
        top=Inches(1.9),
        size=32,
        width=Inches(11.3),
        align=PP_ALIGN.CENTER,
    )
    # URL callout box — gold border, gold-tint fill, matching .demo-url in styles.css
    box_w, box_h = Inches(8.5), Inches(1.1)
    box_left = (SLIDE_W - box_w) / 2
    callout = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, box_left, Inches(2.9), box_w, box_h)
    callout.adjustments[0] = 0.12
    callout.fill.solid()
    callout.fill.fore_color.rgb = GOLD_TINT
    callout.line.color.rgb = GOLD
    callout.line.width = Pt(1.5)
    callout.shadow.inherit = False
    tf = callout.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "lango-app-dusky.vercel.app"
    run.font.name = FONT_MONO
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = GOLD

    box, tf = add_textbox(slide, Inches(1.4), Inches(4.35), Inches(10.5), Inches(0.5))
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Walk through: Command Center → Audit Log → Fairness Audit → Drift & Security → Pilot & Sandbox"
    run.font.name = FONT_MONO
    run.font.size = Pt(15)
    run.font.color.rgb = TEXT_SECONDARY

    box, tf = add_textbox(slide, Inches(2.4), Inches(5.0), Inches(8.5), Inches(1.0))
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = (
        "Live at lango-app-dusky.vercel.app — real Rust/Axum backend on Render, real PostgreSQL, "
        "not a frontend-only mockup. Presenter drives the actual app for this section, not slides."
    )
    run.font.name = FONT_SANS
    run.font.size = Pt(14)
    run.font.color.rgb = TEXT_MUTED
    add_footer(slide)

    # ---- Slide 6 — Fairness evidence (screenshot, real visual) ----
    image_slide(prs, os.path.join(CAPTURES_DIR, "slide-6-capture.png"))

    # ---- Slide 7 — Security/monitoring evidence (screenshot, real visual) ----
    image_slide(prs, os.path.join(CAPTURES_DIR, "slide-7-capture.png"))

    # ---- Slide 8 — Business Model ----
    text_slide(
        prs,
        "08 — Business Model",
        "Institution pays, staff use, compliance benefits",
        [
            "Customer: the institution (bank/hospital/ministry), bought at CISO/Head of Compliance level",
            "Pilot phase: no revenue yet — proving the concept at one institution, one department",
            "Post-pilot: per-seat/per-institution licensing, priced against the cost of a compliance incident avoided",
        ],
    )

    # ---- Slide 9 — Roadmap ----
    text_slide(
        prs,
        "09 — Roadmap",
        "30 / 60 / 90-day path from a real, deployed, hardened v0.1 to a validated pilot",
        [
            "Day 30: pilot institution and department confirmed, consent signed off, institution onboarded onto "
            "the platform (multi-tenant isolation, rate limiting, and a basic internal security pass are already "
            "built and tested — this is about onboarding a real institution onto existing infrastructure, not "
            "building that infrastructure)",
            "Day 60: midpoint review — redaction accuracy and fairness measured on real pilot traffic",
            "Day 90: full pilot cohort onboarded, go/no-go decision on scale-out",
        ],
        headline_size=26,
    )

    # ---- Slide 10 — Team + Ask ----
    text_slide(
        prs,
        "10 — Team + Ask",
        "Team Lango — and what we need next",
        [
            "Phakamile Mlala (Team Leader, Electronic Engineering, NUST Bulawayo) & Vanessa Moyo "
            "(Researcher, Product Design)",
            "Built with Claude and Claude Code for drafting and implementation — reviewed by the team throughout",
            "Ask: a real pilot institution partner — multi-tenant isolation, rate limiting, and a basic internal "
            "security pass are already built and tested, not the ask anymore",
            "Also need: input from a real institutional security/compliance reviewer on the "
            "shared-vs-dedicated-instance tenancy tradeoff (see docs/DEPLOYMENT_PLAN.md) — a decision worth "
            "making with an actual institutional stakeholder, not in isolation — plus support standing up a live "
            "AI provider connection and a formal penetration test ahead of real institutional traffic",
        ],
    )

    out_path = os.path.join(HERE, "Lango_Pitch_Deck.pptx")
    prs.save(out_path)
    print(f"Saved {out_path}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
